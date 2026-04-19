import os
import requests
from bs4 import BeautifulSoup
from graph_api import generate_access_token
import tempfile
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from google import genai
from google.genai import types
try:
    import whisper
except ImportError:
    whisper = None
    print("Warning: whisper not installed — audio/video transcription will be unavailable")
from datetime import datetime, timedelta
import time
import json
from config import Config
from text_cleaning import clean_summary_text

def summarize_content_with_gemini(content):
    client = genai.Client(api_key=Config.GEMINI_API_KEY)
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=['summarize', content],
        config=types.GenerateContentConfig(
        safety_settings=[
            types.SafetySetting(
                category=types.HarmCategory.HARM_CATEGORY_HARASSMENT,
                threshold=types.HarmBlockThreshold.BLOCK_NONE,
            ),
            types.SafetySetting(
                category=types.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                threshold=types.HarmBlockThreshold.BLOCK_NONE,
            ),
            types.SafetySetting(
                category=types.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                threshold=types.HarmBlockThreshold.BLOCK_NONE,
            ),
            types.SafetySetting(
                category=types.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                threshold=types.HarmBlockThreshold.BLOCK_NONE,
            ),
            types.SafetySetting(
                category=types.HarmCategory.HARM_CATEGORY_CIVIC_INTEGRITY,
                threshold=types.HarmBlockThreshold.BLOCK_NONE,
            ),
        ]
        )
    )
    return clean_summary_text(response.text)

GRAPH_API_ENDPOINT = 'https://graph.microsoft.com/v1.0'

def retry_with_backoff(func, max_retries=5, backoff_factor=2, **kwargs):
    """
    Retry a function with exponential backoff.

    Args:
        func (callable): The function to execute.
        max_retries (int): Maximum number of retries.
        backoff_factor (int): Factor by which the wait time increases.
        **kwargs: Arguments to pass to the function.

    Returns:
        The result of the function call, if successful.

    Raises:
        Exception: If all retries fail.
    """
    retries = 0
    while retries < max_retries:
        try:
            return func(**kwargs)
        except requests.RequestException as e:
            retries += 1
            wait_time = backoff_factor ** retries
            print(f"Attempt {retries}/{max_retries} failed. Retrying in {wait_time} seconds...")
            time.sleep(wait_time)
    raise Exception(f"Failed after {max_retries} retries.")


# Lazy-load Whisper model to avoid import-time delay
_whisper_model = None

def _get_whisper_model():
    global _whisper_model
    if _whisper_model is None:
        _whisper_model = whisper.load_model("base")
    return _whisper_model


def fetch_onenote_notebooks(access_token):
    """
    Fetch all OneNote notebooks for the authenticated user.
    """
    headers = {
        "Authorization": f"Bearer {access_token}"
    }

    # Use retry_with_backoff to fetch notebooks
    notebooks_url = "https://graph.microsoft.com/v1.0/me/onenote/notebooks"
    response = retry_with_backoff(requests.get, url=notebooks_url, headers=headers)

    if response.status_code != 200:
        print(f"Failed to fetch notebooks: {response.status_code}, {response.json()}")
        return []

    notebooks = response.json().get('value', [])
    if not notebooks:
        print("No notebooks found.")
        return []

    return notebooks

def fetch_onenote_sections(access_token, notebook_id):
    """
    Fetch sections from a specific OneNote notebook.
    """
    headers = {
        "Authorization": f"Bearer {access_token}"
    }

    # Use retry_with_backoff to fetch sections
    sections_url = f"https://graph.microsoft.com/v1.0/me/onenote/notebooks/{notebook_id}/sections"
    response = retry_with_backoff(requests.get, url=sections_url, headers=headers)

    if response.status_code != 200:
        print(f"Failed to fetch sections: {response.status_code}, {response.json()}")
        return []

    sections = response.json().get('value', [])
    if not sections:
        print(f"No sections found in notebook {notebook_id}.")
        return []

    return sections

def fetch_onenote_content(access_token, target_section_name):
    """
    Fetch OneNote content from all notebooks and their sections using a stack for section navigation.
    """
    target_section_name = os.path.splitext(target_section_name)[0]

    # Fetch all notebooks using retry logic
    notebooks = fetch_onenote_notebooks(access_token)
    if not notebooks:
        return []

    content_list = []

    for notebook in notebooks:
        notebook_id = notebook['id']
        sections = fetch_onenote_sections(access_token, notebook_id)
        if not sections:
            continue

        section_stack = []

        for section in sections:
            section_stack.append(section)

        while section_stack:
            current_section = section_stack.pop()
            section_id = current_section['id']
            section_name = current_section['displayName']

            if section_name != target_section_name:
                continue

            headers = {
                "Authorization": f"Bearer {access_token}"
            }

            # Fetch pages with retry logic
            section_pages_url = f"https://graph.microsoft.com/v1.0/me/onenote/sections/{section_id}/pages"
            response = retry_with_backoff(requests.get, url=section_pages_url, headers=headers)

            if response.status_code != 200:
                print(f"Failed to fetch pages for section {section_name}: {response.status_code}, {response.json()}")
                continue

            pages = response.json().get('value', [])
            if not pages:
                print(f"No pages found in section {section_name}.")
                continue

            for page in pages:
                page_id = page['id']
                page_title = page.get('title', 'Untitled Page')

                # Fetch page content with retry logic
                page_content_url = f"https://graph.microsoft.com/v1.0/me/onenote/pages/{page_id}/content"
                page_response = retry_with_backoff(requests.get, url=page_content_url, headers=headers)

                if page_response.status_code == 200:
                    html_content = page_response.text
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text_content = soup.get_text(separator='\n', strip=True)

                    if text_content.strip():
                        content_list.append({
                            "title": page_title,
                            "content": text_content
                        })
                    else:
                        print(f"No text content found for page: {page_title}")
                else:
                    print(f"Failed to fetch content for page {page_title}: {page_response.status_code}")

    return content_list

def list_onedrive_items(headers, folder_id=None):
    try:
        if folder_id:
            endpoint = f'{GRAPH_API_ENDPOINT}/me/drive/items/{folder_id}/children'
        else:
            endpoint = f'{GRAPH_API_ENDPOINT}/me/drive/root/children'

        response = retry_with_backoff(requests.get, url=endpoint, headers=headers)

        if response.status_code != 200:
            print(f"Error: {response.status_code}")
            print(response.json())
            raise Exception(response.json())

        items = response.json().get('value', [])
        if not items:
            return []

        item_list = []
        for idx, item in enumerate(items):
            name = item.get('name', 'Unnamed')
            item_id = item.get('id', None)
            is_folder = item.get('folder', None) is not None
            item_type_desc = "Folder" if is_folder else "File"

            item_list.append((name, item_id, item_type_desc))
        return item_list
    except Exception as e:
        print(e)
        return []
combined_content = []
processed_files = set()

def get_onedrive_file_content(headers, file_id, file_name, access_token, cutoff_date):
    file_extension = os.path.splitext(file_name)[1].lower()
    url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content"
    file_metadata_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}"

    # Fetch file metadata with retry logic
    response = retry_with_backoff(requests.get, url=file_metadata_url, headers=headers)
    if response.status_code != 200:
        print(f"Failed to fetch file metadata. Status Code: {response.status_code}, Error: {response.text}")
        return "", f"Failed to fetch file metadata (HTTP {response.status_code})"

    file_metadata = response.json()
    last_modified_str = file_metadata.get("lastModifiedDateTime", "Unknown")
    last_modified = datetime.fromisoformat(last_modified_str[:-1])  # Remove 'Z' for parsing

    # Fetch file content with retry logic
    response = retry_with_backoff(requests.get, url=url, headers=headers)
    if response.status_code != 200:
        print(f"Failed to fetch file content. Status Code: {response.status_code}, Error: {response.text}")
        return "", f"Failed to fetch file content (HTTP {response.status_code})"
    
    file_content = response.content
    extracted_text = ""

    # Process audio files using Whisper
    if file_extension in [".mp3", ".wav", ".m4a", ".flac", ".mov"]:
        try:
            model = _get_whisper_model()
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_audio:
                temp_audio.write(file_content)
                temp_audio_path = temp_audio.name
            
            result = model.transcribe(temp_audio_path)
            extracted_text = result["text"]
            os.remove(temp_audio_path)
        except Exception as e:
            print(f"Error transcribing audio file {file_name}: {e}")
            return "", f"Error transcribing audio file: {e}"

    elif file_extension == ".pdf":
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name

            with open(temp_file_path, "rb") as pdf_file:
                reader = PdfReader(pdf_file)
                extracted_text = "\n".join(page.extract_text() or "" for page in reader.pages)

            os.remove(temp_file_path)
        except Exception as e:
            print(f"Error reading PDF file {file_name}: {e}")
            return "", f"Error reading PDF file: {e}"

    elif file_extension == ".docx":
        try:
            doc = Document(BytesIO(file_content))
            extracted_text = "\n".join(para.text for para in doc.paragraphs)
        except Exception as e:
            print(f"Error reading DOCX file {file_name}: {e}")
            return "", f"Error reading Word file: {e}"

    elif file_extension == ".txt":
        try:
            extracted_text = file_content.decode('utf-8')
        except Exception as e:
            print(f"Error reading TXT file {file_name}: {e}")
            return "", f"Error reading text file: {e}"

    elif file_extension in [".xlsx", ".xls"]:
        try:
            wb = openpyxl.load_workbook(BytesIO(file_content), read_only=True, data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()
            extracted_text = "\n\n".join(sheets_text)
        except Exception as e:
            print(f"Error reading Excel file {file_name}: {e}")
            return "", f"Error reading Excel file: {e}"

    elif file_extension == ".csv":
        try:
            extracted_text = file_content.decode('utf-8')
        except Exception as e:
            print(f"Error reading CSV file {file_name}: {e}")
            return "", f"Error reading CSV file: {e}"

    elif file_extension in [".pptx"]:
        try:
            from pptx import Presentation
            prs = Presentation(BytesIO(file_content))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
            extracted_text = "\n\n".join(slides_text)
        except Exception as e:
            print(f"Error reading PPTX file {file_name}: {e}")
            return "", f"Error reading PowerPoint file: {e}"

    elif file_extension == ".one":
        try:
            onenote_content = fetch_onenote_content(access_token['access_token'], file_name)
            extracted_text = "\n".join(page['content'] for page in onenote_content)
        except Exception as e:
            print(f"Error extracting OneNote content from file {file_name}: {e}")
            return "", f"Error reading OneNote file: {e}"

    if not extracted_text.strip():
        return "", "No readable text content found in this file."  # Return empty tuple instead of None

    # Summarize extracted content
    summary = summarize_content_with_gemini(extracted_text)

    # Return structured data
    return extracted_text, summary


def navigate_onedrive(headers, access_token, cutoff_days):
    # Reset state for each new call
    current_folder_id = None
    folder_stack = []
    visited_folders = set()
    processed_files = set()  # Ensure processed files are reset
    file_structure = {}  # Initialize file_structure here
    file_list = []

    # Calculate cutoff date
    cutoff_date = datetime.now() - timedelta(days=cutoff_days)

    while True:
        folder_name = "Root" if not folder_stack else folder_stack[-1][0]
        items = list_onedrive_items(headers, folder_id=current_folder_id)

        if not items:
            if not folder_stack:
                break
            _, current_folder_id = folder_stack.pop()
            continue

        # Store files in the dictionary
        file_structure[folder_name] = [
            item[0] for item in items if item[2] == "File"
        ]

        folder_found = False
        for item in items:
            item_name, item_id, item_type = item

            if item_type == "File" and item_name not in processed_files and not item_name.endswith('.onetoc2'):
                file_metadata_url = f"https://graph.microsoft.com/v1.0/me/drive/items/{item_id}"

                # Fetch file metadata with retry logic
                response = retry_with_backoff(requests.get, url=file_metadata_url, headers=headers)
                if response.status_code != 200:
                    print(f"Failed to fetch file metadata. Status Code: {response.status_code}, Error: {response.text}")
                    return None

                file_metadata = response.json()
                last_modified_str = file_metadata.get("lastModifiedDateTime", "Unknown")
                last_modified = datetime.fromisoformat(last_modified_str[:-1])  # Remove 'Z' for parsing
                print(last_modified, cutoff_date)
                if last_modified >= cutoff_date:
                    processed_files.add(item_name)
                    file_list.append((item_name, item_id))

            elif item_type == "Folder":
                if item_id not in visited_folders:
                    visited_folders.add(item_id)
                    folder_stack.append((folder_name, current_folder_id))
                    current_folder_id = item_id
                    folder_found = True
                    break

        if not folder_found:
            if not folder_stack:
                break
            _, current_folder_id = folder_stack.pop()

    return file_list



def format_combined_content(content_list):
    """
    Dynamically formats a list of dictionaries containing title, last modified, and content.
    """
    formatted_output = "\n\n".join(
        f"Title: {entry.get('title', 'No Title')}\n"
        f"Last Modified: {entry.get('last_modified', 'No Date')}\n"
        f"Content:\n{entry.get('content', 'No Content')}"
        for entry in content_list
    )
    return formatted_output
